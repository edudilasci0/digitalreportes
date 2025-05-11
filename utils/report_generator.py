# utils/report_generator.py

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import io
from datetime import datetime
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
import collections

def generate_excel(metrics, projections, program_analysis, comentarios, marca):
    """Generar informe en formato Excel"""
    buffer = io.BytesIO()
    
    with pd.ExcelWriter(buffer, engine='xlsxwriter') as writer:
        # 1. Hoja de resumen
        metricas_list = [
            'Leads Acumulados',
            'Matrículas Acumuladas',
            'Objetivo de Matrículas',
            'Tasa de Conversión (%)',
            '% Matrículas Leads Nuevos',
            '% Matrículas Remarketing',
            'Inversión Acumulada',
            'CPL Promedio'
        ]
        
        valores_list = [
            metrics['leads_acumulados'],
            metrics['matriculas_acumuladas'],
            metrics['objetivo_matriculas'],
            f"{metrics['tasa_conversion']:.2f}%",
            f"{metrics['pct_matriculas_nuevos']:.1f}%",
            f"{metrics['pct_matriculas_remarketing']:.1f}%",
            f"${metrics['inversion_acumulada']:,.2f}",
            f"${metrics['cpl_promedio']:,.2f}"
        ]
        
        # Agregar tiempo transcurrido solo si es relevante
        if marca in ["GRADO", "UNISUD"] and metrics['tiempo_transcurrido'] is not None:
            metricas_list.insert(0, 'Tiempo Transcurrido (%)')
            valores_list.insert(0, f"{metrics['tiempo_transcurrido']:.1f}%")
        
        df_resumen = pd.DataFrame({
            'Métrica': metricas_list,
            'Valor': valores_list
        })
        
        df_resumen.to_excel(writer, sheet_name='Resumen', index=False)
        
        # 2. Hoja de proyecciones
        df_proyecciones = pd.DataFrame({
            'Métrica': [
                'Leads Proyectados',
                'Matrículas Proyectadas (Min)',
                'Matrículas Proyectadas (Max)',
                '% Cumplimiento Proyectado'
            ],
            'Valor': [
                projections['leads_proyectados'],
                projections['matriculas_proyectadas_min'],
                projections['matriculas_proyectadas_max'],
                f"{projections['pct_cumplimiento_proyectado']:.1f}%"
            ]
        })
        
        df_proyecciones.to_excel(writer, sheet_name='Proyecciones', index=False)
        
        # 3. Hoja de distribución de resultados
        # Asegurarnos que program_analysis['tabla_completa'] sea un DataFrame
        if not isinstance(program_analysis['tabla_completa'], pd.DataFrame):
            print("Convirtiendo 'tabla_completa' a DataFrame")
            df_tabla_completa = pd.DataFrame(program_analysis['tabla_completa'])
        else:
            df_tabla_completa = program_analysis['tabla_completa']
            
        df_tabla_completa.to_excel(writer, sheet_name='Distribución Resultados', index=False)
        
        # 4. Hoja de Top 5 programas
        # Asegurarnos que program_analysis['top_matriculas'] sea un DataFrame
        if not isinstance(program_analysis['top_matriculas'], pd.DataFrame):
            print("Convirtiendo 'top_matriculas' a DataFrame")
            df_top_matriculas = pd.DataFrame(program_analysis['top_matriculas'])
        else:
            df_top_matriculas = program_analysis['top_matriculas']
            
        df_top_matriculas.to_excel(writer, sheet_name='Top Programas', index=False)
        
        # 5. Hoja de programas con menor conversión
        # Asegurarnos que program_analysis['menor_conversion'] sea un DataFrame
        if not isinstance(program_analysis['menor_conversion'], pd.DataFrame):
            print("Convirtiendo 'menor_conversion' a DataFrame")
            df_menor_conversion = pd.DataFrame(program_analysis['menor_conversion'])
        else:
            df_menor_conversion = program_analysis['menor_conversion']
            
        df_menor_conversion.to_excel(writer, sheet_name='Menor Conversión', index=False)
        
        # 6. Hoja de comentarios
        df_comentarios = pd.DataFrame({
            'Fecha': [datetime.now().strftime('%Y-%m-%d %H:%M:%S')],
            'Comentarios': [comentarios]
        })
        
        df_comentarios.to_excel(writer, sheet_name='Comentarios', index=False)
        
        # Dar formato a las hojas
        workbook = writer.book
        
        # Formato para títulos
        header_format = workbook.add_format({
            'bold': True,
            'bg_color': '#4472C4',
            'font_color': 'white',
            'border': 1
        })
        
        # Aplicar formato a cada hoja
        for sheet_name in writer.sheets:
            worksheet = writer.sheets[sheet_name]
            worksheet.set_column('A:A', 30)
            worksheet.set_column('B:Z', 15)
            
            # Aplicar formato a los encabezados de forma segura
            try:
                for col_num, value in enumerate(writer.sheets[sheet_name].table.columns):
                    worksheet.write(0, col_num, value, header_format)
            except AttributeError as e:
                print(f"Error al formatear encabezados en {sheet_name}: {e}")
    
    buffer.seek(0)
    return buffer

def generate_pdf(metrics, projections, program_analysis, comentarios, marca):
    """Generar informe en formato PDF"""
    pdf = FPDF()
    pdf.add_page()
    
    # Configuración de estilo
    pdf.set_font('Arial', 'B', 16)
    pdf.cell(190, 10, f"Reporte Estratégico - {marca}", 0, 1, 'C')
    pdf.ln(10)
    
    # 1. Estado Actual
    pdf.set_font('Arial', 'B', 14)
    pdf.cell(190, 10, "Estado Actual", 0, 1, 'L')
    pdf.ln(5)
    
    pdf.set_font('Arial', '', 12)
    
    # Solo mostrar tiempo transcurrido para marcas con convocatorias
    if marca in ["GRADO", "UNISUD"] and metrics['tiempo_transcurrido'] is not None:
        pdf.cell(95, 10, f"Tiempo Transcurrido: {metrics['tiempo_transcurrido']:.1f}%", 0, 0, 'L')
        pdf.cell(95, 10, f"Leads Acumulados: {metrics['leads_acumulados']}", 0, 1, 'L')
    else:
        pdf.cell(190, 10, f"Leads Acumulados: {metrics['leads_acumulados']}", 0, 1, 'L')
        
    pdf.cell(95, 10, f"Matrículas vs Objetivo: {metrics['matriculas_acumuladas']}/{metrics['objetivo_matriculas']}", 0, 0, 'L')
    pdf.cell(95, 10, f"Tasa de Conversión: {metrics['tasa_conversion']:.2f}%", 0, 1, 'L')
    pdf.ln(10)
    
    # 2. Composición de Resultados
    pdf.set_font('Arial', 'B', 14)
    pdf.cell(190, 10, "Composición de Resultados", 0, 1, 'L')
    pdf.ln(5)
    
    pdf.set_font('Arial', '', 12)
    pdf.cell(95, 10, f"% Matrículas Leads Nuevos: {metrics['pct_matriculas_nuevos']:.1f}%", 0, 0, 'L')
    pdf.cell(95, 10, f"% Matrículas Remarketing: {metrics['pct_matriculas_remarketing']:.1f}%", 0, 1, 'L')
    pdf.ln(10)
    
    # 3. Estimación de Cierre
    pdf.set_font('Arial', 'B', 14)
    pdf.cell(190, 10, "Estimación de Cierre", 0, 1, 'L')
    pdf.ln(5)
    
    pdf.set_font('Arial', '', 12)
    pdf.cell(95, 10, f"Leads Proyectados: {projections['leads_proyectados']}", 0, 0, 'L')
    pdf.cell(95, 10, f"Matrículas Proyectadas: {projections['matriculas_proyectadas_min']} - {projections['matriculas_proyectadas_max']}", 0, 1, 'L')
    pdf.cell(190, 10, f"% Cumplimiento Proyectado: {projections['pct_cumplimiento_proyectado']:.1f}%", 0, 1, 'L')
    pdf.ln(10)
    
    # 4. Top 5 Programas
    pdf.set_font('Arial', 'B', 14)
    pdf.cell(190, 10, "Top 5 Programas con Más Matrículas", 0, 1, 'L')
    pdf.ln(5)
    
    # Crear tabla de top 5 programas
    pdf.set_font('Arial', 'B', 10)
    pdf.cell(95, 10, "Programa", 1, 0, 'C')
    pdf.cell(30, 10, "Leads", 1, 0, 'C')
    pdf.cell(30, 10, "Matrículas", 1, 0, 'C')
    pdf.cell(35, 10, "Tasa Conv. (%)", 1, 1, 'C')
    
    # Asegurarnos que program_analysis['top_matriculas'] sea un DataFrame
    if not isinstance(program_analysis['top_matriculas'], pd.DataFrame):
        df_top_matriculas = pd.DataFrame(program_analysis['top_matriculas'])
    else:
        df_top_matriculas = program_analysis['top_matriculas']
    
    pdf.set_font('Arial', '', 10)
    for _, row in df_top_matriculas.iterrows():
        pdf.cell(95, 10, str(row['Programa']), 1, 0, 'L')
        pdf.cell(30, 10, str(row['Leads']), 1, 0, 'C')
        pdf.cell(30, 10, str(row['Matrículas']), 1, 0, 'C')
        pdf.cell(35, 10, str(row['Tasa Conversión (%)']), 1, 1, 'C')
    
    pdf.ln(10)
    
    # 5. Comentarios
    pdf.set_font('Arial', 'B', 14)
    pdf.cell(190, 10, "Comentarios", 0, 1, 'L')
    pdf.ln(5)
    
    pdf.set_font('Arial', '', 12)
    pdf.multi_cell(190, 10, comentarios)
    
    # Convertir a buffer
    buffer = io.BytesIO()
    buffer.write(pdf.output(dest='S').encode('latin1'))
    buffer.seek(0)
    
    return buffer

def generate_pptx(metrics, projections, program_analysis, comentarios, marca):
    """Generar presentación en formato PowerPoint en formato horizontal"""
    prs = Presentation()
    
    # Cambiar orientación a horizontal (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 1. Portada
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = f"Reporte Status Semanal - {marca}"
    subtitle = slide.placeholders[1]
    subtitle.text = f"Fecha: {datetime.now().strftime('%Y-%m-%d')}"
    
    # 2. Estado Actual
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Estado Actual"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    # Solo mostrar tiempo transcurrido para marcas con convocatorias
    if marca in ["GRADO", "UNISUD"] and metrics['tiempo_transcurrido'] is not None:
        tf.text = f"Tiempo Transcurrido: {metrics['tiempo_transcurrido']:.1f}%\n"
    else:
        tf.text = ""
        
    tf.text += f"Leads Acumulados: {metrics['leads_acumulados']}\n"
    tf.text += f"Matrículas vs Objetivo: {metrics['matriculas_acumuladas']}/{metrics['objetivo_matriculas']}\n"
    tf.text += f"Tasa de Conversión: {metrics['tasa_conversion']:.2f}%\n"
    tf.text += f"Proyección de Cumplimiento: {projections['pct_cumplimiento_proyectado']:.1f}%"
    
    # Agregar gráfico de barras para el progreso
    if marca in ["GRADO", "UNISUD"] and metrics['tiempo_transcurrido'] is not None:
        left = Inches(7)
        top = Inches(2)
        width = Inches(5)
        height = Inches(0.5)
        
        # Agregar rectángulo de fondo
        shape = slide.shapes.add_shape(
            1, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = (225, 225, 225)  # Gris claro
        shape.line.color.rgb = (200, 200, 200)       # Borde gris
        
        # Agregar rectángulo de progreso
        progress_width = width * (metrics['tiempo_transcurrido'] / 100)
        progress_shape = slide.shapes.add_shape(
            1, left, top, progress_width, height
        )
        progress_shape.fill.solid()
        progress_shape.fill.fore_color.rgb = (0, 112, 192)  # Azul
        progress_shape.line.fill.background()  # Sin borde
        
        # Texto encima de la barra
        text_box = slide.shapes.add_textbox(left, top - Inches(0.3), width, Inches(0.25))
        text_box.text = "Tiempo Transcurrido"
        
        # Texto con el porcentaje
        text_box = slide.shapes.add_textbox(left + width / 2 - Inches(0.5), top + Inches(0.15), Inches(1), Inches(0.25))
        text_frame = text_box.text_frame
        text_frame.text = f"{metrics['tiempo_transcurrido']:.1f}%"
        
        # Barra para matrículas vs objetivo
        top = top + Inches(1.2)
        
        # Agregar rectángulo de fondo
        shape = slide.shapes.add_shape(
            1, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = (225, 225, 225)  # Gris claro
        shape.line.color.rgb = (200, 200, 200)       # Borde gris
        
        # Agregar rectángulo de progreso
        pct_objetivo = min(1.0, metrics['matriculas_acumuladas'] / max(1, metrics['objetivo_matriculas']))
        progress_width = width * pct_objetivo
        progress_shape = slide.shapes.add_shape(
            1, left, top, progress_width, height
        )
        progress_shape.fill.solid()
        progress_shape.fill.fore_color.rgb = (112, 173, 71)  # Verde
        progress_shape.line.fill.background()  # Sin borde
        
        # Texto encima de la barra
        text_box = slide.shapes.add_textbox(left, top - Inches(0.3), width, Inches(0.25))
        text_box.text = "Matrículas vs Objetivo"
        
        # Texto con el porcentaje
        text_box = slide.shapes.add_textbox(left + width / 2 - Inches(0.5), top + Inches(0.15), Inches(1), Inches(0.25))
        text_frame = text_box.text_frame
        text_frame.text = f"{metrics['matriculas_acumuladas']}/{metrics['objetivo_matriculas']}"
    
    # 3. Composición de Resultados
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Composición de Resultados"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Distribución de Matrículas por tipo de Lead:\n\n"
    tf.text += f"Leads Nuevos: {metrics['pct_matriculas_nuevos']:.1f}%\n"
    tf.text += f"Remarketing: {metrics['pct_matriculas_remarketing']:.1f}%"
    
    # Agregar gráfico de barras para la composición
    left = Inches(7)
    top = Inches(2)
    width = Inches(5)
    height = Inches(0.5)
    
    # Barra para leads nuevos
    shape = slide.shapes.add_shape(
        1, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = (225, 225, 225)  # Gris claro
    shape.line.color.rgb = (200, 200, 200)       # Borde gris
    
    progress_width = width * (metrics['pct_matriculas_nuevos'] / 100)
    progress_shape = slide.shapes.add_shape(
        1, left, top, progress_width, height
    )
    progress_shape.fill.solid()
    progress_shape.fill.fore_color.rgb = (0, 112, 192)  # Azul
    progress_shape.line.fill.background()  # Sin borde
    
    text_box = slide.shapes.add_textbox(left, top - Inches(0.3), width, Inches(0.25))
    text_box.text = "Leads Nuevos"
    
    text_box = slide.shapes.add_textbox(left + progress_width + Inches(0.1), top + Inches(0.15), Inches(1), Inches(0.25))
    text_frame = text_box.text_frame
    text_frame.text = f"{metrics['pct_matriculas_nuevos']:.1f}%"
    
    # Barra para remarketing
    top = top + Inches(1.2)
    
    shape = slide.shapes.add_shape(
        1, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = (225, 225, 225)  # Gris claro
    shape.line.color.rgb = (200, 200, 200)       # Borde gris
    
    progress_width = width * (metrics['pct_matriculas_remarketing'] / 100)
    progress_shape = slide.shapes.add_shape(
        1, left, top, progress_width, height
    )
    progress_shape.fill.solid()
    progress_shape.fill.fore_color.rgb = (255, 153, 0)  # Naranja
    progress_shape.line.fill.background()  # Sin borde
    
    text_box = slide.shapes.add_textbox(left, top - Inches(0.3), width, Inches(0.25))
    text_box.text = "Remarketing"
    
    text_box = slide.shapes.add_textbox(left + progress_width + Inches(0.1), top + Inches(0.15), Inches(1), Inches(0.25))
    text_frame = text_box.text_frame
    text_frame.text = f"{metrics['pct_matriculas_remarketing']:.1f}%"
    
    # 4. Estimación de Cierre
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Estimación de Cierre"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = f"Leads Proyectados: {projections['leads_proyectados']} ± {int(projections['leads_proyectados_std'])}\n"
    tf.text += f"Matrículas Proyectadas: {projections['matriculas_proyectadas_mean']} ± {int(projections['matriculas_proyectadas_std'])}\n"
    tf.text += f"Intervalo 90% Confianza: {projections['matriculas_proyectadas_min']} - {projections['matriculas_proyectadas_max']}\n\n"
    tf.text += f"Probabilidades de Alcanzar Objetivo:"
    
    # Agregar barras para las probabilidades
    left = Inches(1)
    top = Inches(3.2)
    width = Inches(4)
    height = Inches(0.3)
    
    umbrales = [80, 90, 100, 110, 120]
    for i, umbral in enumerate(umbrales):
        prob_key = f'prob_meta_{umbral}'
        curr_top = top + Inches(i * 0.5)
        
        # Marco de fondo
        shape = slide.shapes.add_shape(
            1, left, curr_top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = (225, 225, 225)  # Gris claro
        shape.line.color.rgb = (200, 200, 200)       # Borde gris
        
        # Barra de progreso
        prob_pct = min(1.0, projections[prob_key] / 100)
        progress_width = width * prob_pct
        progress_shape = slide.shapes.add_shape(
            1, left, curr_top, progress_width, height
        )
        progress_shape.fill.solid()
        
        # Color según probabilidad
        if projections[prob_key] >= 75:
            color_rgb = (112, 173, 71)  # Verde
        elif projections[prob_key] >= 50:
            color_rgb = (255, 192, 0)   # Amarillo
        else:
            color_rgb = (237, 125, 49)  # Naranja/Rojo
            
        progress_shape.fill.fore_color.rgb = color_rgb
        progress_shape.line.fill.background()  # Sin borde
        
        # Etiqueta
        text_box = slide.shapes.add_textbox(left - Inches(1.5), curr_top, Inches(1.4), height)
        text_frame = text_box.text_frame
        text_frame.text = f"{umbral}% del Objetivo"
        
        # Valor
        text_box = slide.shapes.add_textbox(left + width + Inches(0.1), curr_top, Inches(1.5), height)
        text_frame = text_box.text_frame
        text_frame.text = f"{projections[prob_key]:.1f}% prob."
    
    # Agregar imagen de la distribución a la derecha
    # (Esto requeriría crear y guardar una imagen temporalmente, lo cual es complejo)
    # En su lugar, añadimos un texto explicativo
    text_box = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(5), Inches(4))
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = "Distribución de Matrículas Proyectadas"
    p.font.bold = True
    p.font.size = Pt(14)
    
    p = text_frame.add_paragraph()
    p.text = f"La simulación Monte Carlo muestra que con un nivel de confianza del 90%, se espera obtener entre {projections['matriculas_proyectadas_min']} y {projections['matriculas_proyectadas_max']} matrículas al cierre de la convocatoria."
    
    p = text_frame.add_paragraph()
    p.text = f"Si el objetivo es {metrics['objetivo_matriculas']} matrículas, la probabilidad de alcanzarlo es del {projections['prob_meta_100']:.1f}%."
    
    # 5. Top 5 Programas
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Top 5 Programas con Más Matrículas"
    
    # Asegurarnos que program_analysis['top_matriculas'] sea un DataFrame
    if not isinstance(program_analysis['top_matriculas'], pd.DataFrame):
        df_top_matriculas = pd.DataFrame(program_analysis['top_matriculas'])
    else:
        df_top_matriculas = program_analysis['top_matriculas']
    
    # Crear tabla
    if not df_top_matriculas.empty:
        rows = len(df_top_matriculas) + 1  # +1 para el encabezado
        cols = 4
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(11)
        height = Inches(0.5 * rows)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Ajustar anchos de columna
        table.columns[0].width = Inches(6)  # Programa (más ancho)
        table.columns[1].width = Inches(1.5)  # Leads
        table.columns[2].width = Inches(1.5)  # Matrículas
        table.columns[3].width = Inches(2)    # Tasa Conv.
        
        # Encabezados
        table.cell(0, 0).text = "Programa"
        table.cell(0, 1).text = "Leads"
        table.cell(0, 2).text = "Matrículas"
        table.cell(0, 3).text = "Tasa Conv. (%)"
        
        # Dar formato a los encabezados
        for i in range(cols):
            cell = table.cell(0, i)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (0, 112, 192)  # Azul
            cell.text_frame.paragraphs[0].font.color.rgb = (255, 255, 255)  # Texto blanco
            cell.text_frame.paragraphs[0].font.bold = True
        
        # Datos
        for i, (_, row) in enumerate(df_top_matriculas.iterrows(), 1):
            table.cell(i, 0).text = str(row['Programa'])
            table.cell(i, 1).text = str(row['Leads'])
            table.cell(i, 2).text = str(row['Matrículas'])
            table.cell(i, 3).text = str(row['Tasa Conversión (%)'])
            
            # Alternar colores de fila
            if i % 2 == 0:
                for j in range(cols):
                    cell = table.cell(i, j)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = (240, 240, 240)  # Gris muy claro
    else:
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "No hay datos disponibles para mostrar"
    
    # 6. Programas con Menor Conversión
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Programas con Menor Conversión"
    
    # Asegurarnos que program_analysis['menor_conversion'] sea un DataFrame
    if not isinstance(program_analysis['menor_conversion'], pd.DataFrame):
        df_menor_conversion = pd.DataFrame(program_analysis['menor_conversion'])
    else:
        df_menor_conversion = program_analysis['menor_conversion']
    
    # Crear tabla similar al slide anterior
    if not df_menor_conversion.empty:
        rows = len(df_menor_conversion) + 1  # +1 para el encabezado
        cols = 4
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(11)
        height = Inches(0.5 * rows)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Ajustar anchos de columna
        table.columns[0].width = Inches(6)  # Programa (más ancho)
        table.columns[1].width = Inches(1.5)  # Leads
        table.columns[2].width = Inches(1.5)  # Matrículas
        table.columns[3].width = Inches(2)    # Tasa Conv.
        
        # Encabezados
        table.cell(0, 0).text = "Programa"
        table.cell(0, 1).text = "Leads"
        table.cell(0, 2).text = "Matrículas"
        table.cell(0, 3).text = "Tasa Conv. (%)"
        
        # Dar formato a los encabezados
        for i in range(cols):
            cell = table.cell(0, i)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (192, 0, 0)  # Rojo
            cell.text_frame.paragraphs[0].font.color.rgb = (255, 255, 255)  # Texto blanco
            cell.text_frame.paragraphs[0].font.bold = True
        
        # Datos
        for i, (_, row) in enumerate(df_menor_conversion.iterrows(), 1):
            table.cell(i, 0).text = str(row['Programa'])
            table.cell(i, 1).text = str(row['Leads'])
            table.cell(i, 2).text = str(row['Matrículas'])
            table.cell(i, 3).text = str(row['Tasa Conversión (%)'])
            
            # Alternar colores de fila
            if i % 2 == 0:
                for j in range(cols):
                    cell = table.cell(i, j)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = (240, 240, 240)  # Gris muy claro
    else:
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "No hay datos disponibles para mostrar"
    
    # 7. Comentarios
    if comentarios and comentarios.strip():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Comentarios"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = comentarios
    
    # Guardar en buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    return buffer 